#!/usr/bin/env python3
"""生成对外汇报用 Word / PowerPoint（领导可直接打开）。"""
from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt

HERE = Path(__file__).resolve().parent
OUT = HERE.parent / "docs" / "platform-briefing" / "对外汇报"


def add_doc_title(doc: Document, text: str) -> None:
    p = doc.add_paragraph()
    r = p.add_run(text)
    r.bold = True
    r.font.size = Pt(22)
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER


def build_docx() -> None:
    doc = Document()
    add_doc_title(doc, "AI Agent 供应链数字化平台")
    sub = doc.add_paragraph("对外汇报材料（简版）")
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph("文档版本：V1.0").alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph("说明：本平台为自主研发的一体化「AI + 数据智能」能力，业务经营库采用只读方式对接，不替代核心 ERP 下单与主数据维护。").alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_page_break()

    doc.add_heading("一、项目定位", level=1)
    doc.add_paragraph(
        "面向集团供应链与数字化创新场景，提供统一门户：数据驾驶舱、天枢大屏、数据洞察、"
        "AI 训练与识别、票据智能处理、智能物流等能力，实现「看得见、算得清、可演示、可审计」的运营视图。"
    )

    doc.add_heading("二、建设目标", level=1)
    for t in (
        "数字化创新：大屏与驾驶舱集中呈现经营与履约态势；",
        "AI 赋能供应链：图像识别、模型训练部署、单据识别与比对服务稽核与效率；",
        "自主研发平台：技术栈与交付形态自主可控，支持私有化与集团合规要求。",
    ):
        doc.add_paragraph(t, style="List Bullet")

    doc.add_heading("三、系统边界（只读解耦）", level=1)
    doc.add_paragraph(
        "业务主数据与订单写入仍在既有业务系统中完成；本平台通过只读账号访问业务 MySQL，"
        "分析接口仅允许 SELECT；平台自有 SQLite 承载数据集、训练任务、模型与物流绑定等平台侧数据。"
    )
    table = doc.add_table(rows=4, cols=2)
    table.style = "Table Grid"
    rows = [
        ("业务 MySQL（洞察）", "只读查询与聚合，支撑驾驶舱、洞察 API、实时经营通道"),
        ("平台 SQLite", "数据集、训练、模型元数据、物流相关等平台数据"),
        ("物流数据源", "按租户切库，与既有后台逻辑对齐"),
        ("核心业务写入", "不在本平台内完成下单、新增商品、报价、收货等写操作"),
    ]
    for i, (a, b) in enumerate(rows):
        table.rows[i].cells[0].text = a
        table.rows[i].cells[1].text = b

    doc.add_heading("四、总体架构示意", level=1)
    arch = OUT / "platform-architecture-mindmap.png"
    if arch.exists():
        doc.add_picture(str(arch), width=Inches(6.2))
    doc.add_paragraph("上图：分层架构意象示意（展示层 / 应用层 / 数据层）。")

    doc.add_heading("五、业务只读与平台关系", level=1)
    rel = OUT / "对外汇报-只读洞察与平台关系图.png"
    if rel.exists():
        doc.add_picture(str(rel), width=Inches(6.2))

    doc.add_heading("六、核心能力摘要", level=1)
    caps = [
        "数据洞察：KPI、分布、地图、热力、运营告警等；支持多表名/主键自适应，降低对接成本。",
        "实时经营通道：在零 DDL 前提下基于增量轮询 + WebSocket 推送今日 GMV/单量等，供驾驶舱与天枢「实时链路」展示。",
        "AI 训练与推理：数据集管理、训练任务、模型库与部署、在线识别 Top-K。",
        "票据智能：图像识别与结构化输出，支持与订单信息比对，服务稽核场景。",
        "智能物流：车辆、多源设备绑定、轨迹与地图、费用管理；坐标系与地图服务可配置。",
    ]
    for c in caps:
        doc.add_paragraph(c, style="List Bullet")

    doc.add_heading("七、AI 训练—部署—推理闭环", level=1)
    ai_png = OUT / "对外汇报-AI训练推理闭环.png"
    if ai_png.exists():
        doc.add_picture(str(ai_png), width=Inches(6.2))

    doc.add_heading("八、典型应用场景", level=1)
    for t in (
        "管理端：日常运营、模型与数据管理、洞察分析、物流管理；",
        "指挥参观：驾驶舱 / 天枢大屏集中演示数字化与实时经营态势；",
        "技术支撑：只读 SQL 分析、开放 API 与集团集成。",
    ):
        doc.add_paragraph(t, style="List Bullet")

    doc.add_heading("九、对外表述关键词", level=1)
    for t in (
        "自主研发的一体化 AI + 数据智能平台；",
        "业务库只读解耦，降低对生产交易系统风险；",
        "AI 赋能供应链（训练/推理/单据）；",
        "实时增量洞察，零 DDL 适配运维约束。",
    ):
        doc.add_paragraph(t, style="List Bullet")

    doc.add_heading("十、说明", level=1)
    doc.add_paragraph(
        "本材料不含环境账号口令；正式对外版本可由信息中心替换为集团标准模板页眉页脚。"
        "更细的技术说明见同目录上级 Markdown 与研发归档。"
    )

    out = OUT / "对外汇报-供应链数字化平台.docx"
    doc.save(str(out))
    print("Wrote", out)


def add_slide_title(prs: Presentation, title: str, subtitle: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.4), PptxInches(9), PptxInches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PptxPt(28)
    p.font.bold = True
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = PptxPt(14)


def add_slide_bullets(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.35), PptxInches(9), PptxInches(0.9))
    tb.text_frame.paragraphs[0].text = title
    tb.text_frame.paragraphs[0].font.size = PptxPt(24)
    tb.text_frame.paragraphs[0].font.bold = True
    body = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(1.2), PptxInches(9), PptxInches(5.5))
    tf = body.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        para = tf.add_paragraph() if i else tf.paragraphs[0]
        if i:
            para = tf.add_paragraph()
        para.text = b
        para.font.size = PptxPt(16)
        para.level = 0


def add_slide_picture(prs: Presentation, title: str, png: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(9), PptxInches(0.8))
    tb.text_frame.paragraphs[0].text = title
    tb.text_frame.paragraphs[0].font.size = PptxPt(22)
    tb.text_frame.paragraphs[0].font.bold = True
    if png.exists():
        slide.shapes.add_picture(str(png), PptxInches(0.5), PptxInches(1.1), width=PptxInches(9))


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    # 封面
    s0 = prs.slides.add_slide(prs.slide_layouts[6])
    t = s0.shapes.add_textbox(PptxInches(1), PptxInches(2.2), PptxInches(11), PptxInches(2))
    t.text_frame.paragraphs[0].text = "AI Agent 供应链数字化平台"
    t.text_frame.paragraphs[0].font.size = PptxPt(36)
    t.text_frame.paragraphs[0].font.bold = True
    st = s0.shapes.add_textbox(PptxInches(1), PptxInches(3.2), PptxInches(11), PptxInches(1))
    st.text_frame.paragraphs[0].text = "对外汇报｜数字化转型 · AI 赋能 · 自主研发"
    st.text_frame.paragraphs[0].font.size = PptxPt(18)
    foot = s0.shapes.add_textbox(PptxInches(1), PptxInches(6.2), PptxInches(11), PptxInches(0.8))
    foot.text_frame.paragraphs[0].text = "材料版本 V1.0（不含敏感连接信息）"
    foot.text_frame.paragraphs[0].font.size = PptxPt(12)

    add_slide_bullets(
        prs,
        "项目定位",
        [
            "自主研发的一体化「AI + 数据智能」平台",
            "服务供应链可视化、稽核与运营指挥场景",
            "与既有业务系统协同：经营库只读，不替代 ERP 核心写入",
        ],
    )
    add_slide_bullets(
        prs,
        "建设目标",
        [
            "数字化创新：驾驶舱、天枢大屏、数据洞察",
            "AI 赋能：训练 / 部署 / 识别 / 票据智能",
            "可交付：私有化、合规叙事清晰、可持续迭代",
        ],
    )
    add_slide_bullets(
        prs,
        "系统边界（领导关切）",
        [
            "业务 MySQL：只读账号，SELECT 聚合与分析",
            "平台 SQLite：数据集、训练、模型、物流绑定等平台数据",
            "不在此平台完成：下单、主数据维护、收货入账等核心写操作",
        ],
    )
    add_slide_picture(prs, "总体架构示意", OUT / "platform-architecture-mindmap.png")
    add_slide_picture(prs, "只读洞察与平台关系", OUT / "对外汇报-只读洞察与平台关系图.png")
    add_slide_bullets(
        prs,
        "核心能力（一）数据与实时",
        [
            "多维度洞察：KPI、区域/品类、地图、热力、告警等",
            "自适应表结构：订单/明细表名与主键差异自动适配，降低对接成本",
            "实时经营通道：增量轮询 + WebSocket，零 DDL 展示今日 GMV/单量",
        ],
    )
    add_slide_picture(prs, "AI 训练—部署—推理闭环", OUT / "对外汇报-AI训练推理闭环.png")
    add_slide_bullets(
        prs,
        "核心能力（二）AI 与票据",
        [
            "数据集与训练任务管理，模型库与一键部署",
            "在线识别服务，可对接验收与品控场景",
            "票据识别与结构化比对，服务稽核与效率提升",
        ],
    )
    add_slide_bullets(
        prs,
        "核心能力（三）智能物流",
        [
            "车辆与多源设备（定位/视频/温湿度等）绑定",
            "轨迹与地图展示，坐标系可配置",
            "费用录入与统计",
        ],
    )
    add_slide_bullets(
        prs,
        "典型场景",
        [
            "管理端：日常运营、模型与数据、洞察与物流",
            "指挥参观：驾驶舱 / 天枢大屏演示",
            "技术支撑：只读分析与 API 集成",
        ],
    )
    add_slide_bullets(
        prs,
        "对外表述关键词",
        [
            "自主研发、数字化创新、AI 赋能供应链",
            "业务库只读解耦、风险可控",
            "实时增量洞察、零 DDL",
        ],
    )
    end = prs.slides.add_slide(prs.slide_layouts[6])
    et = end.shapes.add_textbox(PptxInches(2), PptxInches(3), PptxInches(9), PptxInches(1.5))
    et.text_frame.paragraphs[0].text = "谢谢聆听 · 欢迎指导"
    et.text_frame.paragraphs[0].font.size = PptxPt(32)

    outp = OUT / "对外汇报-供应链数字化平台.pptx"
    prs.save(str(outp))
    print("Wrote", outp)


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    build_docx()
    build_pptx()


if __name__ == "__main__":
    main()
