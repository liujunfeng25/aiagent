"""Markdown → .docx 转换。

用 markdown-it-py 的 Token 流扫描（足以覆盖我们约定的：
  H1/H2/H3、段落、加粗、无序/有序列表、表格），
再用 python-docx 写入 Word 文档。

不追求全量 Markdown，主要覆盖 System Prompt 要求的报告结构；
解析失败时退化为「每行一段落」，保证不抛异常。
"""

from __future__ import annotations

import io
import logging
import re
from typing import Iterable

logger = logging.getLogger(__name__)


def markdown_to_docx_bytes(markdown: str, *, title: str | None = None) -> bytes:
    """返回一个 .docx 文件的二进制字节。"""
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Microsoft YaHei"
    style.font.size = Pt(11)

    if title:
        doc.add_heading(title, level=0)

    try:
        _render_with_markdown_it(doc, markdown)
    except Exception as e:
        logger.warning("markdown-it 解析失败，按段落回退：%s", e)
        _render_fallback(doc, markdown)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# markdown-it-py 渲染
# ---------------------------------------------------------------------------

def _render_with_markdown_it(doc, markdown: str) -> None:
    from markdown_it import MarkdownIt

    md = MarkdownIt("commonmark", {"html": False}).enable("table")
    tokens = md.parse(markdown or "")

    i = 0
    n = len(tokens)
    while i < n:
        t = tokens[i]
        typ = t.type

        if typ == "heading_open":
            level = int(t.tag.lstrip("h") or 1)
            inline = tokens[i + 1] if i + 1 < n else None
            text = _inline_text(inline)
            doc.add_heading(text, level=min(max(level, 1), 3))
            i += 3
            continue

        if typ == "paragraph_open":
            inline = tokens[i + 1] if i + 1 < n else None
            p = doc.add_paragraph()
            _write_inline_with_runs(p, inline)
            i += 3
            continue

        if typ == "bullet_list_open":
            i = _render_list(doc, tokens, i, ordered=False)
            continue

        if typ == "ordered_list_open":
            i = _render_list(doc, tokens, i, ordered=True)
            continue

        if typ == "table_open":
            i = _render_table(doc, tokens, i)
            continue

        if typ in ("hr", "softbreak"):
            doc.add_paragraph("—" * 10)
            i += 1
            continue

        if typ == "fence" or typ == "code_block":
            p = doc.add_paragraph()
            r = p.add_run(t.content or "")
            r.font.name = "Courier New"
            i += 1
            continue

        i += 1


def _inline_text(inline) -> str:
    if inline is None or inline.type != "inline":
        return ""
    return "".join(
        c.content for c in (inline.children or []) if getattr(c, "type", "") == "text"
    ) or (inline.content or "")


def _write_inline_with_runs(paragraph, inline) -> None:
    """把 inline 的 children（text/strong/em/code/link/softbreak）写成带格式的 run。"""
    if inline is None or inline.type != "inline":
        if inline is not None and inline.content:
            paragraph.add_run(inline.content)
        return

    bold = False
    italic = False
    for c in inline.children or []:
        t = getattr(c, "type", "")
        if t == "strong_open":
            bold = True
        elif t == "strong_close":
            bold = False
        elif t == "em_open":
            italic = True
        elif t == "em_close":
            italic = False
        elif t == "text":
            r = paragraph.add_run(c.content)
            r.bold = bold
            r.italic = italic
        elif t == "code_inline":
            r = paragraph.add_run(c.content)
            r.font.name = "Courier New"
        elif t == "softbreak" or t == "hardbreak":
            paragraph.add_run("\n")
        elif t == "link_open":
            # 简化：把链接文本当普通 run 输出
            pass
        elif t == "link_close":
            pass


def _render_list(doc, tokens, i, *, ordered: bool) -> int:
    """处理 bullet_list_open...bullet_list_close 区间，返回下一个 i。"""
    style_name = "List Number" if ordered else "List Bullet"
    # 跳过当前 list_open
    i += 1
    depth = 1
    while i < len(tokens) and depth > 0:
        t = tokens[i]
        if t.type in ("bullet_list_open", "ordered_list_open"):
            depth += 1
            i += 1
            continue
        if t.type in ("bullet_list_close", "ordered_list_close"):
            depth -= 1
            i += 1
            continue
        if t.type == "list_item_open":
            # 下一条应当是 paragraph_open/inline/paragraph_close
            j = i + 1
            text_parts: list = []
            while j < len(tokens) and tokens[j].type != "list_item_close":
                if tokens[j].type == "inline":
                    text_parts.append(tokens[j])
                j += 1
            try:
                p = doc.add_paragraph(style=style_name)
            except KeyError:
                p = doc.add_paragraph()
                p.add_run("• " if not ordered else "1. ")
            for inl in text_parts:
                _write_inline_with_runs(p, inl)
            i = j + 1
            continue
        i += 1
    return i


def _render_table(doc, tokens, i) -> int:
    """处理 table_open...table_close；返回下一个 i。"""
    rows: list[list[list]] = []  # [row][col] = inline tokens list
    header_rows = 0
    i += 1
    while i < len(tokens) and tokens[i].type != "table_close":
        t = tokens[i]
        if t.type in ("thead_open", "tbody_open"):
            is_header = t.type == "thead_open"
            i += 1
            while i < len(tokens) and tokens[i].type not in ("thead_close", "tbody_close"):
                if tokens[i].type == "tr_open":
                    row: list[list] = []
                    i += 1
                    while i < len(tokens) and tokens[i].type != "tr_close":
                        if tokens[i].type in ("th_open", "td_open"):
                            j = i + 1
                            inline = None
                            while j < len(tokens) and tokens[j].type not in ("th_close", "td_close"):
                                if tokens[j].type == "inline":
                                    inline = tokens[j]
                                j += 1
                            row.append([inline] if inline else [])
                            i = j + 1
                            continue
                        i += 1
                    rows.append(row)
                    if is_header:
                        header_rows = max(header_rows, 1)
                    i += 1
                    continue
                i += 1
            i += 1
            continue
        i += 1

    if rows:
        col_count = max((len(r) for r in rows), default=0) or 1
        table = doc.add_table(rows=len(rows), cols=col_count)
        try:
            table.style = "Light Grid Accent 1"
        except KeyError:
            pass
        for ri, row in enumerate(rows):
            for ci in range(col_count):
                cell = table.cell(ri, ci)
                cell.text = ""
                if ci < len(row):
                    for inl in row[ci]:
                        p = cell.paragraphs[0] if cell.paragraphs else cell.add_paragraph()
                        if p.runs:
                            p = cell.add_paragraph()
                        _write_inline_with_runs(p, inl)
                if ri < header_rows and cell.paragraphs:
                    for run in cell.paragraphs[0].runs:
                        run.bold = True
    # 跳过 table_close
    return i + 1


# ---------------------------------------------------------------------------
# 回退：按行生成段落
# ---------------------------------------------------------------------------

_HEAD_RE = re.compile(r"^(#{1,3})\s+(.*)")


def _render_fallback(doc, markdown: str) -> None:
    for raw in (markdown or "").splitlines():
        line = raw.rstrip()
        if not line.strip():
            doc.add_paragraph("")
            continue
        m = _HEAD_RE.match(line)
        if m:
            doc.add_heading(m.group(2).strip(), level=len(m.group(1)))
            continue
        doc.add_paragraph(line)


# ---------------------------------------------------------------------------
# Markdown → .md （原文保底导出）
# ---------------------------------------------------------------------------

def markdown_to_md_bytes(markdown: str, *, title: str | None = None) -> bytes:
    """直接导出原始 Markdown；若给了 title 且正文未以 `# ` 开头，补一个一级标题。"""
    text = (markdown or "").strip()
    if title and not re.match(r"^#\s+", text):
        text = f"# {title}\n\n{text}"
    return text.encode("utf-8")


# ---------------------------------------------------------------------------
# Markdown → .pptx
#
# 设计：
#   - 一个 `# 标题` → 一张封面页（如已有 title 参数，则封面用 title）
#   - 一个 `## 小节` → 一张内容页；页内的段落 / 列表 / 表格依次注入正文框
#   - 表格单独成页（整段复用 python-pptx 的 add_table），跟随最近的 ##
# 不追求漂亮排版，目标是给汇报准备一个可以二次编辑的版本。
# ---------------------------------------------------------------------------

def markdown_to_pptx_bytes(markdown: str, *, title: str | None = None) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from markdown_it import MarkdownIt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 封面
    cover_layout = prs.slide_layouts[0]
    cover = prs.slides.add_slide(cover_layout)
    try:
        cover.shapes.title.text = title or "业务报告"
        if len(cover.placeholders) > 1:
            cover.placeholders[1].text = "食迅易联 · AI 业务分析助手自动生成"
    except Exception:
        pass

    md = MarkdownIt("commonmark", {"html": False}).enable("table")
    tokens = md.parse(markdown or "")

    content_layout = prs.slide_layouts[5]  # Title Only
    current_slide = None
    current_body_lines: list[str] = []

    def _flush_body():
        """把累积的正文行写到当前 slide 里。"""
        nonlocal current_slide, current_body_lines
        if current_slide is None or not current_body_lines:
            current_body_lines = []
            return
        left = Inches(0.6)
        top = Inches(1.4)
        width = Inches(12.1)
        height = Inches(5.7)
        tx = current_slide.shapes.add_textbox(left, top, width, height)
        tf = tx.text_frame
        tf.word_wrap = True
        for idx, line in enumerate(current_body_lines):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.name = "Microsoft YaHei"
        current_body_lines = []

    def _new_content_slide(title_text: str) -> None:
        nonlocal current_slide
        _flush_body()
        current_slide = prs.slides.add_slide(content_layout)
        try:
            current_slide.shapes.title.text = title_text
        except Exception:
            pass

    def _table_slide(rows: list[list[str]], header_rows: int) -> None:
        nonlocal current_slide
        _flush_body()
        slide = prs.slides.add_slide(content_layout)
        try:
            slide.shapes.title.text = ""
        except Exception:
            pass
        if not rows:
            return
        nrow = len(rows)
        ncol = max((len(r) for r in rows), default=1) or 1
        left = Inches(0.6)
        top = Inches(1.4)
        width = Inches(12.1)
        height = Inches(5.5)
        table = slide.shapes.add_table(nrow, ncol, left, top, width, height).table
        for ri, row in enumerate(rows):
            for ci in range(ncol):
                cell = table.cell(ri, ci)
                cell.text = row[ci] if ci < len(row) else ""
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(13)
                        r.font.name = "Microsoft YaHei"
                        if ri < header_rows:
                            r.font.bold = True

    i = 0
    n = len(tokens)
    while i < n:
        t = tokens[i]
        typ = t.type

        if typ == "heading_open":
            level = int(t.tag.lstrip("h") or 1)
            inline = tokens[i + 1] if i + 1 < n else None
            text = _inline_text(inline)
            if level <= 1:
                if current_slide is not None:
                    _flush_body()
                cover2 = prs.slides.add_slide(cover_layout)
                try:
                    cover2.shapes.title.text = text
                except Exception:
                    pass
                current_slide = None
            else:
                _new_content_slide(text)
            i += 3
            continue

        if typ == "paragraph_open":
            inline = tokens[i + 1] if i + 1 < n else None
            text = _inline_text(inline)
            if text.strip():
                if current_slide is None:
                    _new_content_slide(title or "报告")
                current_body_lines.append(text.strip())
            i += 3
            continue

        if typ in ("bullet_list_open", "ordered_list_open"):
            if current_slide is None:
                _new_content_slide(title or "报告")
            ordered = typ == "ordered_list_open"
            items, i = _collect_list_items(tokens, i, ordered=ordered)
            for idx, it in enumerate(items):
                bullet = f"{idx + 1}. " if ordered else "• "
                current_body_lines.append(f"{bullet}{it}")
            continue

        if typ == "table_open":
            rows, header_rows, i = _collect_table(tokens, i)
            _table_slide(rows, header_rows)
            continue

        if typ == "fence" or typ == "code_block":
            if current_slide is None:
                _new_content_slide(title or "代码块")
            for ln in (t.content or "").splitlines():
                current_body_lines.append(ln)
            i += 1
            continue

        i += 1

    _flush_body()

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _collect_list_items(tokens, i, *, ordered: bool) -> tuple[list[str], int]:
    """把 list_open..list_close 区间里每个 item 压成一行纯文本。"""
    items: list[str] = []
    i += 1
    depth = 1
    while i < len(tokens) and depth > 0:
        t = tokens[i]
        if t.type in ("bullet_list_open", "ordered_list_open"):
            depth += 1
            i += 1
            continue
        if t.type in ("bullet_list_close", "ordered_list_close"):
            depth -= 1
            i += 1
            continue
        if t.type == "list_item_open":
            j = i + 1
            parts: list[str] = []
            while j < len(tokens) and tokens[j].type != "list_item_close":
                if tokens[j].type == "inline":
                    parts.append(_inline_text(tokens[j]))
                j += 1
            items.append(" ".join(p for p in parts if p).strip())
            i = j + 1
            continue
        i += 1
    return items, i


def _collect_table(tokens, i) -> tuple[list[list[str]], int, int]:
    """把 Markdown 表格 flatten 成字符串行，返回 (rows, header_rows, next_i)。"""
    rows: list[list[str]] = []
    header_rows = 0
    i += 1
    while i < len(tokens) and tokens[i].type != "table_close":
        t = tokens[i]
        if t.type in ("thead_open", "tbody_open"):
            is_header = t.type == "thead_open"
            i += 1
            while i < len(tokens) and tokens[i].type not in ("thead_close", "tbody_close"):
                if tokens[i].type == "tr_open":
                    row: list[str] = []
                    i += 1
                    while i < len(tokens) and tokens[i].type != "tr_close":
                        if tokens[i].type in ("th_open", "td_open"):
                            j = i + 1
                            text = ""
                            while j < len(tokens) and tokens[j].type not in ("th_close", "td_close"):
                                if tokens[j].type == "inline":
                                    text = _inline_text(tokens[j])
                                j += 1
                            row.append(text)
                            i = j + 1
                            continue
                        i += 1
                    rows.append(row)
                    if is_header:
                        header_rows = max(header_rows, 1)
                    i += 1
                    continue
                i += 1
            i += 1
            continue
        i += 1
    return rows, header_rows, i + 1
